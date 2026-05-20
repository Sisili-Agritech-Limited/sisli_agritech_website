from pptx import Presentation

def extract_text(filepath):
    try:
        prs = Presentation(filepath)
        text_content = []
        for i, slide in enumerate(prs.slides):
            text_content.append(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
                
                # Also check notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_content.append(f"NOTES: {notes}")
        return "\n".join(text_content)
    except Exception as e:
        return f"Error reading {filepath}: {str(e)}"

files = [
    "Sisili_Agritech PITCH WITH NOTES.pptx",
    "sisili_agritech_pitch jim igite pitch [Autosaved].pptx"
]

with open("extracted_pitch_text.txt", "w", encoding="utf-8") as out_file:
    for f in files:
        out_file.write(f"=== CONTENT FOR {f} ===\n")
        out_file.write(extract_text(f))
        out_file.write("\n\n")

print("Text extraction complete. Saved to extracted_pitch_text.txt")
